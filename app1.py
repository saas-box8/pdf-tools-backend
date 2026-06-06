"""
PowerPoint (PPT/PPTX) → PDF converter — pure Python, no LibreOffice.
Uses python-pptx to read slides and reportlab to write the PDF.
Works on Render free plan.
"""

from __future__ import annotations

import io
import os
from pathlib import Path

from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib import colors as rl_colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

app = Flask(__name__)
CORS(app)
app.config["MAX_CONTENT_LENGTH"] = int(os.getenv("MAX_UPLOAD_MB", "100")) * 1024 * 1024

ALLOWED_EXTENSIONS = {".ppt", ".pptx"}


def allowed_file(filename: str) -> bool:
    return Path(filename.lower()).suffix in ALLOWED_EXTENSIONS


# ── Alignment map ─────────────────────────────────────────────────────────────
_ALIGN = {
    PP_ALIGN.LEFT: TA_LEFT,
    PP_ALIGN.CENTER: TA_CENTER,
    PP_ALIGN.RIGHT: TA_RIGHT,
    None: TA_LEFT,
}


def _escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _run_markup(run) -> str:
    """Convert a pptx run to reportlab XML markup."""
    text = _escape(run.text)
    if not text:
        return ""
    parts = []
    close = []
    font = run.font
    sz = font.size.pt if font.size else 18
    clr = None
    try:
        if font.color and font.color.rgb:
            r, g, b = font.color.rgb.red, font.color.rgb.green, font.color.rgb.blue
            clr = "#{:02x}{:02x}{:02x}".format(r, g, b)
    except Exception:
        pass
    fc = f'<font size="{sz:.1f}"'
    if clr:
        fc += f' color="{clr}"'
    fc += ">"
    parts.append(fc); close.insert(0, "</font>")
    if font.bold:
        parts.append("<b>"); close.insert(0, "</b>")
    if font.italic:
        parts.append("<i>"); close.insert(0, "</i>")
    if font.underline:
        parts.append("<u>"); close.insert(0, "</u>")
    return "".join(parts) + text + "".join(close)


def pptx_to_pdf_bytes(pptx_bytes: bytes) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))

    # Determine page size from presentation dimensions
    slide_w = prs.slide_width.inches if prs.slide_width else 10
    slide_h = prs.slide_height.inches if prs.slide_height else 7.5
    page_w = slide_w * inch
    page_h = slide_h * inch

    buf = io.BytesIO()
    pdf = SimpleDocTemplate(
        buf,
        pagesize=(page_w, page_h),
        leftMargin=0.5 * inch,
        rightMargin=0.5 * inch,
        topMargin=0.5 * inch,
        bottomMargin=0.5 * inch,
    )

    story = []
    title_style = ParagraphStyle("title", fontSize=24, leading=28, bold=True,
                                  alignment=TA_CENTER, spaceAfter=12,
                                  textColor=rl_colors.HexColor("#1E3A5F"))
    body_style = ParagraphStyle("body", fontSize=14, leading=18, spaceAfter=6)
    bullet_style = ParagraphStyle("bullet", fontSize=13, leading=17,
                                   leftIndent=18, spaceAfter=4,
                                   bulletIndent=6)
    slide_num_style = ParagraphStyle("num", fontSize=8, leading=10,
                                      alignment=TA_RIGHT,
                                      textColor=rl_colors.HexColor("#888888"))

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx > 0:
            story.append(PageBreak())

        story.append(Paragraph(f"Slide {slide_idx + 1}", slide_num_style))
        story.append(Spacer(1, 4))

        # Sort shapes top-to-bottom by vertical position
        shapes = sorted(
            slide.shapes,
            key=lambda s: (s.top if s.top is not None else 0)
        )

        for shape in shapes:
            if not shape.has_text_frame:
                continue

            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                raw_text = para.text.strip()
                if not raw_text:
                    story.append(Spacer(1, 3))
                    continue

                # Build markup from runs
                markup = "".join(_run_markup(r) for r in para.runs)
                if not markup.strip():
                    markup = _escape(raw_text)

                align = _ALIGN.get(para.alignment, TA_LEFT)

                # Detect title placeholder
                is_title = (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                    and para_idx == 0
                )

                if is_title:
                    ps = ParagraphStyle("slide_title", parent=title_style, alignment=align)
                else:
                    level = para.level or 0
                    indent = level * 12
                    ps = ParagraphStyle(
                        f"p_{slide_idx}_{para_idx}",
                        parent=body_style,
                        alignment=align,
                        leftIndent=indent,
                        fontSize=max(9, 14 - level * 1.5),
                        leading=max(11, 18 - level * 2),
                    )

                try:
                    story.append(Paragraph(markup, ps))
                except Exception:
                    story.append(Paragraph(_escape(raw_text), ps))

    if not story:
        story.append(Paragraph("(empty presentation)", body_style))

    pdf.build(story)
    buf.seek(0)
    return buf.read()


@app.route("/", methods=["GET"])
def home():
    return jsonify({"status": "running", "tool": "powerpoint to pdf"})


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"ok": True}), 200


@app.route("/convert", methods=["POST"])
def convert():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        file = request.files["file"]
        if not file.filename:
            return jsonify({"error": "Empty filename"}), 400
        if not allowed_file(file.filename):
            return jsonify({"error": "Only PPT or PPTX files are allowed"}), 400

        safe_name = secure_filename(file.filename)
        pptx_bytes = file.read()
        pdf_bytes = pptx_to_pdf_bytes(pptx_bytes)

        response = send_file(
            io.BytesIO(pdf_bytes),
            as_attachment=True,
            download_name=f"{Path(safe_name).stem}.pdf",
            mimetype="application/pdf",
            conditional=False,
            max_age=0,
        )
        response.headers["Cache-Control"] = "no-store"
        return response

    except Exception as e:
        return jsonify({"error": "Conversion failed", "details": str(e)}), 500


if __name__ == "__main__":
    port = int(os.getenv("PORT", "5000"))
    app.run(host="0.0.0.0", port=port, debug=False)
