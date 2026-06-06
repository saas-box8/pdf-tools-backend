import io
import math
import os
import logging

import fitz
import cv2
import numpy as np

from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from werkzeug.utils import secure_filename

# ─── App Setup ────────────────────────────────────────────────────────────────

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = int(os.getenv("MAX_UPLOAD_MB", "100")) * 1024 * 1024
app.config["JSON_SORT_KEYS"] = False

cors_origins = os.getenv("CORS_ORIGINS", "*").strip()
if cors_origins == "*":
    CORS(app)
else:
    allowed = [o.strip() for o in cors_origins.split(",") if o.strip()]
    CORS(app, resources={r"/*": {"origins": allowed}})

logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO"))
logger = logging.getLogger(__name__)

EMU_PER_PT = 12700

# ── PyMuPDF text-span flag bitmasks ──────────────────────────────────────────
# (from PyMuPDF docs – these were SWAPPED in the original code, which is a
#  critical bug: bold was read as italic and italic as superscript)
PDF_FLAG_SUPERSCRIPT = 1    # bit 0
PDF_FLAG_ITALIC      = 2    # bit 1
PDF_FLAG_SERIFED     = 4    # bit 2  (font classification, not a style)
PDF_FLAG_MONOSPACE   = 8    # bit 3  (font classification, not a style)
PDF_FLAG_BOLD        = 16   # bit 4


# ─── Basic Helpers ────────────────────────────────────────────────────────────

def points_to_inches(v): return v / 72.0
def pt_to_emu(v): return int(round(v * EMU_PER_PT))
def clamp(n, lo, hi): return max(lo, min(n, hi))


def parse_float(value, default, min_value=None, max_value=None):
    try:
        parsed = float(value)
    except Exception:
        parsed = default
    if min_value is not None:
        parsed = max(min_value, parsed)
    if max_value is not None:
        parsed = min(max_value, parsed)
    return parsed


# ─── Color ────────────────────────────────────────────────────────────────────

def pdf_color_to_rgb(color_value):
    """Convert a PyMuPDF color value (int or float-tuple) to pptx RGBColor."""
    if color_value is None:
        return RGBColor(0, 0, 0)
    try:
        if isinstance(color_value, int):
            return RGBColor(
                (color_value >> 16) & 255,
                (color_value >> 8)  & 255,
                 color_value        & 255,
            )
        if isinstance(color_value, (tuple, list)) and len(color_value) >= 3:
            r, g, b = color_value[:3]
            if any(isinstance(c, float) for c in (r, g, b)):
                r = int(clamp(round(r * 255), 0, 255))
                g = int(clamp(round(g * 255), 0, 255))
                b = int(clamp(round(b * 255), 0, 255))
            else:
                r = int(clamp(r, 0, 255))
                g = int(clamp(g, 0, 255))
                b = int(clamp(b, 0, 255))
            return RGBColor(r, g, b)
    except Exception:
        pass
    return RGBColor(0, 0, 0)


# ─── Font Helpers ─────────────────────────────────────────────────────────────

# Style suffixes that python-pptx controls separately via bold/italic flags
_FONT_SUFFIXES = (
    "-BoldItalicMT", "-BoldItalic", "-BoldMT", "-Bold",
    "-ItalicMT", "-Italic", ",BoldItalic", ",Bold", ",Italic", "MT",
)


def normalize_font_name(font_name):
    """Clean up a PDF embedded font name to a sensible family name."""
    if not font_name:
        return "Arial"
    # Strip subset prefix (e.g. "ABCDEF+Arial-Bold" → "Arial-Bold")
    if "+" in font_name:
        font_name = font_name.split("+", 1)[1]
    # Strip style suffixes so pptx bold/italic attributes take over
    for suffix in _FONT_SUFFIXES:
        if font_name.endswith(suffix):
            font_name = font_name[: -len(suffix)]
            break
    return font_name.strip() or "Arial"


# ─── Slide Background ─────────────────────────────────────────────────────────

def get_page_bg_color(page):
    """
    Sample the page's dominant background colour by rendering it at a tiny
    size and taking the median of corner/edge pixels.
    """
    mat = fitz.Matrix(0.15, 0.15)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    w, h = img.size
    if w < 1 or h < 1:
        return (255, 255, 255)

    pts = [
        (0, 0), (max(0, w-1), 0), (0, max(0, h-1)), (max(0, w-1), max(0, h-1)),
        (w//4, 0), (3*w//4, 0),
        (w//4, max(0, h-1)), (3*w//4, max(0, h-1)),
        (0, h//4), (0, 3*h//4),
        (max(0, w-1), h//4), (max(0, w-1), 3*h//4),
    ]
    samples = [img.getpixel((px, py)) for px, py in pts if 0 <= px < w and 0 <= py < h]
    if not samples:
        return (255, 255, 255)

    r = sorted(s[0] for s in samples)[len(samples)//2]
    g = sorted(s[1] for s in samples)[len(samples)//2]
    b = sorted(s[2] for s in samples)[len(samples)//2]
    return (r, g, b)


def apply_slide_background(slide, rgb):
    """Set a solid background colour on a slide (skip near-white)."""
    r, g, b = rgb
    if r >= 248 and g >= 248 and b >= 248:
        return  # default white – nothing to do
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


# ─── Drawing Analysis (underline / strikethrough) ─────────────────────────────

def get_underline_rects(page):
    """
    Return bounding boxes of thin horizontal lines drawn on the page.
    These correspond to underlines or strikethroughs.
    """
    rects = []
    try:
        for path in page.get_drawings():
            r = path.get("rect")
            if r is None:
                continue
            # Underlines are very thin (height ≤ 2 pt) and reasonably wide
            if 0 < r.height <= 2.0 and r.width > 4:
                rects.append((r.x0, r.y0, r.x1, r.y1))
    except Exception:
        pass
    return rects


def is_span_underlined(span_bbox, underline_rects, tol=3.0):
    """True if a horizontal line sits just below this span."""
    if not underline_rects or not span_bbox:
        return False
    x0, y0, x1, y1 = span_bbox
    for ux0, uy0, ux1, uy1 in underline_rects:
        if uy0 >= y1 - tol and uy1 <= y1 + tol * 2:   # below text baseline
            if ux0 < x1 and ux1 > x0:                  # horizontal overlap
                return True
    return False


def is_span_struck(span_bbox, line_rects, tol=3.0):
    """True if a horizontal line crosses through the vertical midpoint of this span."""
    if not line_rects or not span_bbox:
        return False
    x0, y0, x1, y1 = span_bbox
    mid_y = (y0 + y1) / 2
    for rx0, ry0, rx1, ry1 in line_rects:
        if ry0 >= y0 + tol and ry1 <= y1 - tol:        # inside the span vertically
            if abs((ry0 + ry1)/2 - mid_y) < tol * 2:  # near the midpoint
                if rx0 < x1 and rx1 > x0:
                    return True
    return False


# ─── Text Alignment Detection ─────────────────────────────────────────────────

def detect_line_alignment(line, block):
    """
    Guess the paragraph alignment of a text line by comparing its bounding
    box to its containing block.
    """
    lb = line.get("bbox")
    bb = block.get("bbox")
    if not lb or not bb:
        return PP_ALIGN.LEFT

    lx0, _, lx1, _ = lb
    bx0, _, bx1, _ = bb
    block_w = bx1 - bx0
    if block_w < 1:
        return PP_ALIGN.LEFT

    left_gap  = lx0 - bx0
    right_gap = bx1 - lx1
    line_w    = lx1 - lx0

    # Centre: gaps on both sides are roughly equal and non-trivial
    if abs(left_gap - right_gap) < max(4, block_w * 0.04) and left_gap > 3:
        return PP_ALIGN.CENTER

    # Right-aligned: nearly flush right, significant left indent
    if right_gap < 3 and left_gap > block_w * 0.15:
        return PP_ALIGN.RIGHT

    # Justified: line fills ≥ 90 % of block width (skip the last line heuristic)
    if line_w >= block_w * 0.90 and block_w > 50:
        return PP_ALIGN.JUSTIFY

    return PP_ALIGN.LEFT


# ─── Hyperlinks ───────────────────────────────────────────────────────────────

def get_page_links(page):
    """Return [{uri, rect}] for all URI hyperlinks on the page."""
    links = []
    try:
        for link in page.get_links():
            if link.get("kind") == fitz.LINK_URI and link.get("uri"):
                links.append({"uri": link["uri"], "rect": link["from"]})
    except Exception:
        pass
    return links


def find_link_for_span(span_bbox, page_links, tol=2.0):
    """Return the URI that covers this span, or None."""
    if not span_bbox or not page_links:
        return None
    x0, y0, x1, y1 = span_bbox
    for link in page_links:
        r = link["rect"]
        if not (x1 < r.x0 - tol or x0 > r.x1 + tol or y1 < r.y0 - tol or y0 > r.y1 + tol):
            return link["uri"]
    return None


# ─── Superscript / Baseline Shift via XML ─────────────────────────────────────

def apply_baseline_shift(run, baseline_pct: int):
    """
    Set the <a:rPr baseline="…"> attribute on a run.
    Positive values → superscript, negative → subscript.
    """
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn("baseline"), str(baseline_pct))
    except Exception:
        pass


# ─── Background Rendering (text-erased raster) ────────────────────────────────

def render_background(page, zoom=3.0):
    """
    Render the full page at high DPI, then inpaint (erase) all text bounding
    boxes so that only graphics, images, and decorative elements remain.
    The result is placed behind the editable text layer.
    """
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False, annots=False)
    pil = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    img = np.array(pil)

    mask = np.zeros((img.shape[0], img.shape[1]), dtype=np.uint8)
    sx = pix.width  / page.rect.width
    sy = pix.height / page.rect.height

    for w in page.get_text("words") or []:
        l = clamp(int(w[0] * sx) - 2, 0, pix.width  - 1)
        t = clamp(int(w[1] * sy) - 2, 0, pix.height - 1)
        r = clamp(int(w[2] * sx) + 2, 0, pix.width  - 1)
        b = clamp(int(w[3] * sy) + 2, 0, pix.height - 1)
        if r > l and b > t:
            cv2.rectangle(mask, (l, t), (r, b), 255, -1)

    mask = cv2.dilate(mask, np.ones((3, 3), np.uint8), iterations=1)
    inpainted = cv2.inpaint(img, mask, 5, cv2.INPAINT_TELEA)
    return Image.fromarray(inpainted)


# ─── Core: Add One Text Line as a Textbox ─────────────────────────────────────

def add_line_as_textbox(slide, line, block, scale, offset_x, offset_y,
                         underline_rects, page_links):
    """
    Place a single PDF text line as a precisely-positioned, fully-styled
    textbox on the PowerPoint slide.
    """
    bbox = line.get("bbox")
    if not bbox:
        return

    x0, y0, x1, y1 = bbox
    line_h = max(1.0, y1 - y0)

    # Small horizontal padding prevents text clipping at box edges
    pad_x  = line_h * 0.15
    left   = offset_x + pt_to_emu((x0 - pad_x) * scale)
    top    = offset_y + pt_to_emu(y0 * scale)
    width  = max(1, pt_to_emu((x1 - x0 + 2 * pad_x) * scale))
    # Slight height buffer so descenders aren't cropped
    height = max(1, pt_to_emu(line_h * scale * 1.25))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap   = False
    tf.auto_size   = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment    = detect_line_alignment(line, block)
    p.space_before = Pt(0)
    p.space_after  = Pt(0)
    p.line_spacing = 1.0

    # ── Rotation ──────────────────────────────────────────────────────────
    direction = line.get("dir", (1, 0))
    if isinstance(direction, (tuple, list)) and len(direction) == 2:
        dx, dy = direction
        if not (abs(dx - 1.0) < 1e-6 and abs(dy) < 1e-6):
            try:
                textbox.rotation = -math.degrees(math.atan2(dy, dx))
            except Exception:
                pass

    # ── Spans ─────────────────────────────────────────────────────────────
    for span in line.get("spans", []):
        text = span.get("text", "")
        if not text:
            continue

        run = p.add_run()
        run.text = text

        # — Font family —
        run.font.name = normalize_font_name(span.get("font", ""))

        # — Font size —
        size = span.get("size", 12)
        run.font.size = Pt(max(1, size * scale))

        # — Bold / Italic (CORRECTED – original had these swapped) —
        flags = span.get("flags", 0)
        run.font.bold   = bool(flags & PDF_FLAG_BOLD)
        run.font.italic = bool(flags & PDF_FLAG_ITALIC)

        # — Underline (detected from drawn paths) —
        span_bbox = span.get("bbox")
        if is_span_underlined(span_bbox, underline_rects):
            run.font.underline = True

        # — Strikethrough (drawn paths through vertical midpoint) —
        if is_span_struck(span_bbox, underline_rects):
            try:
                rPr = run._r.get_or_add_rPr()
                rPr.set(qn("strike"), "sngStrike")
            except Exception:
                pass

        # — Superscript —
        if flags & PDF_FLAG_SUPERSCRIPT:
            apply_baseline_shift(run, 30000)   # +30 % above baseline

        # — Font colour —
        color = span.get("color")
        if color is not None:
            run.font.color.rgb = pdf_color_to_rgb(color)

        # — Hyperlinks —
        uri = find_link_for_span(span_bbox, page_links)
        if uri:
            try:
                run.hyperlink.address = uri
            except Exception:
                pass

        # — Character spacing (origin-based, for tightly-tracked fonts) —
        # PyMuPDF provides per-char origins in rawdict; we approximate by
        # checking whether the span's advance differs from its bbox width.
        # If the PDF used letter-spacing, widen/narrow via XML.
        try:
            span_w  = (span_bbox[2] - span_bbox[0]) if span_bbox else 0
            char_ct = len(text)
            if char_ct > 1 and size > 0 and span_w > 0:
                expected_w = size * char_ct * 0.5   # very rough estimate
                ratio      = span_w / max(expected_w, 1)
                if ratio > 1.35:                    # expanded
                    spc_pt = int((ratio - 1.0) * size * 50)
                    rPr = run._r.get_or_add_rPr()
                    rPr.set(qn("spc"), str(spc_pt))
                elif ratio < 0.70:                  # compressed
                    spc_pt = int((ratio - 1.0) * size * 50)
                    rPr = run._r.get_or_add_rPr()
                    rPr.set(qn("spc"), str(max(spc_pt, -200)))
        except Exception:
            pass


# ─── Per-page Slide Population ────────────────────────────────────────────────

def add_page_to_slide(slide, page, ppt, mode="editable", page_zoom=3.0):
    slide_w_emu = ppt.slide_width
    slide_h_emu = ppt.slide_height
    page_w_pt   = page.rect.width
    page_h_pt   = page.rect.height

    scale = min(
        slide_w_emu / pt_to_emu(page_w_pt),
        slide_h_emu / pt_to_emu(page_h_pt),
    )
    page_w_fit = pt_to_emu(page_w_pt * scale)
    page_h_fit = pt_to_emu(page_h_pt * scale)
    offset_x   = int((slide_w_emu - page_w_fit) / 2)
    offset_y   = int((slide_h_emu - page_h_fit) / 2)

    # ── REPLICA: full raster only, no editable layer ───────────────────────
    if mode == "replica":
        pix = page.get_pixmap(
            matrix=fitz.Matrix(max(page_zoom, 3.0), max(page_zoom, 3.0)),
            alpha=False, annots=False,
        )
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        buf = io.BytesIO(); img.save(buf, "PNG"); buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, width=slide_w_emu, height=slide_h_emu)
        return

    # ── Slide background colour ────────────────────────────────────────────
    apply_slide_background(slide, get_page_bg_color(page))

    # ── Background image layer (graphics, photos, decorations – no text) ──
    bg_img = render_background(page, zoom=page_zoom)
    buf = io.BytesIO(); bg_img.save(buf, "PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, offset_x, offset_y, width=page_w_fit, height=page_h_fit)

    # ── Editable text layer ────────────────────────────────────────────────
    underline_rects = get_underline_rects(page)
    page_links      = get_page_links(page)
    text_dict       = page.get_text("dict", sort=True) or {}

    for block in text_dict.get("blocks", []):
        if block.get("type") != 0:        # 0 = text, 1 = image
            continue
        for line in block.get("lines", []):
            if not line.get("spans"):
                continue
            try:
                add_line_as_textbox(
                    slide, line, block, scale, offset_x, offset_y,
                    underline_rects, page_links,
                )
            except Exception as exc:
                logger.warning("Skipped a text line: %s", exc)


# ─── Routes ───────────────────────────────────────────────────────────────────

@app.route("/", methods=["GET"])
def home():
    return jsonify({
        "status": "running",
        "mode": "high-fidelity editable pdf → pptx",
    })


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"ok": True}), 200


@app.route("/convert", methods=["POST"])
def convert():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file      = request.files["file"]
        pdf_bytes = file.read()
        if not pdf_bytes:
            return jsonify({"error": "Empty file uploaded"}), 400

        quality = parse_float(request.form.get("quality", "2"), 2.0, 0.5, 8.0)
        mode    = request.form.get("conversionMode", "editable").strip().lower()
        if mode not in ("editable", "balanced", "replica"):
            mode = "editable"

        if mode == "replica":
            page_zoom = max(quality, 3.0)
        elif mode == "balanced":
            page_zoom = clamp(quality, 1.5, 2.5)
        else:
            page_zoom = max(quality, 2.5)

        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        if doc.page_count == 0:
            return jsonify({"error": "PDF has no pages"}), 400

        first_rect = doc.load_page(0).rect
        ppt = Presentation()
        ppt.slide_width  = Inches(points_to_inches(first_rect.width))
        ppt.slide_height = Inches(points_to_inches(first_rect.height))
        blank_layout = ppt.slide_layouts[6]

        for page in doc:
            slide = ppt.slides.add_slide(blank_layout)
            add_page_to_slide(slide, page, ppt, mode=mode, page_zoom=page_zoom)

        output = io.BytesIO()
        ppt.save(output)
        output.seek(0)

        base_name = secure_filename(file.filename or "converted.pdf")
        root_name = os.path.splitext(base_name)[0] or "converted"

        return send_file(
            output,
            as_attachment=True,
            download_name=f"{root_name}.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            conditional=False,
        )

    except Exception:
        logger.exception("Conversion failed")
        return jsonify({"error": "Conversion failed"}), 500


@app.errorhandler(413)
def file_too_large(_):
    return jsonify({"error": "File too large"}), 413


# ─── Local launcher ───────────────────────────────────────────────────────────

if __name__ == "__main__":
    port  = int(os.getenv("PORT", "5000"))
    debug = os.getenv("FLASK_DEBUG", "0") == "1"
    app.run(host="0.0.0.0", port=port, debug=debug)